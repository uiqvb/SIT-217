# -*- coding: utf-8 -*-
# Creates a 5-slide PPT with two charts for the DronePad summary

import matplotlib
matplotlib.use("Agg")  # headless backend
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime

# ---- personalise here ----
YOUR_NAME = "Your Name"
STUDENT_ID = "sXXXXXXX"
# (Optional) replace with your real pilot numbers:
UTILISATION = 78        # %
AVG_WAIT_MIN = 3.2      # minutes
CONFLICT_PREVENT = 99.2 # %
RES_SUCCESS = 99.85     # %
# --------------------------

ppt_path = "DronePad_5_Slide_Summary.pptx"
kpi_chart_path = "kpi_chart.png"
trend_chart_path = "reservation_trend.png"

# Charts
labels = ["Utilisation %", "Avg wait (min)", "Conflict prevention %", "Reservation success %"]
values = [UTILISATION, AVG_WAIT_MIN, CONFLICT_PREVENT, RES_SUCCESS]

plt.figure()
plt.bar(labels, values)
plt.title("Prototype KPIs (pilot assumptions)")
plt.ylabel("Value")
plt.xticks(rotation=20, ha="right")
plt.tight_layout()
plt.savefig(kpi_chart_path, dpi=200)
plt.close()

hours = np.arange(8, 20)  # 08..19
reservations = np.array([2,3,4,6,7,6,5,4,6,7,5,3])  # demo
capacity = np.array([5]*len(hours))
plt.figure()
plt.plot(hours, reservations, marker="o", label="Reservations/hour")
plt.plot(hours, capacity, marker="s", linestyle="--", label="Capacity/hour")
plt.title("Demand vs Capacity (sample day)")
plt.xlabel("Hour of day")
plt.ylabel("Count")
plt.legend()
plt.tight_layout()
plt.savefig(trend_chart_path, dpi=200)
plt.close()

# PPT helpers
prs = Presentation()
today = datetime.now().strftime("%d %b %Y")

def add_title(title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = title
    s.placeholders[1].text = subtitle

def add_bullets(title, bullets, size=22):
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = title
    tf = s.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        p.text = b
        for r in p.runs:
            r.font.size = Pt(size)

def add_images(title, imgs):
    s = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    s.shapes.title.text = title
    left = Inches(0.5); top = Inches(1.5); width = Inches(4.5)
    for i, img in enumerate(imgs):
        s.shapes.add_picture(img, left + Inches(4.7)*i, top, width=width)

# Slide 1 — Title
add_title(
    "Autonomous Drone Delivery Pad Scheduling – Summary",
    f"SIT217 • Deakin University (Burwood Campus)\n{YOUR_NAME} ({STUDENT_ID})\n{today}"
)

# Slide 2 — Problem, Vision & Scope
add_bullets("Problem, Vision & Scope", [
    "Problem: manual/first-come pad use → clashes, delays, safety risks; weather/airspace inconsistently applied.",
    "Vision: scheduling service allocating pads/time slots, enforcing separation, and gating on weather/airspace; operator control.",
    "In scope: search & reserve, reassignment, check-in/auto-release, notifications, operator console, audit.",
    "Out of scope: flight control algorithms, onboard DAA, hardware procurement."
])

# Slide 3 — High-Level Requirements
add_bullets("High-Level Requirements (FR / QR / BR)", [
    "FR: search & reservation; separation per pad/payload; check-in & auto-release; reassignment; notifications; audit.",
    "QR: P95 search ≤ 1.8 s; ≥99.8% success; 99.95% uptime; ≥99.99% conflict detection; security/privacy; WCAG 2.1 AA.",
    "BR: go-live before T3 pilot; 12-month logs; chargeback; KPI tracking (utilisation, wait time, conflicts avoided)."
], size=20)

# Slide 4 — Graphs
add_images("Prototype Metrics (Illustrative)", [kpi_chart_path, trend_chart_path])

# Slide 5 — V&V, Risks & Next Steps
add_bullets("V&V, Risks & Next Steps", [
    "V&V: Gherkin acceptance; large-scale conflict simulations; performance & security tests; pilot on N1/H2/S1 pads.",
    "Risks: sensor misreads → overrides; sudden weather → dynamic reassign; multi-fleet contention → priority/quotas; API abuse → rate limits.",
    "Next: integrate live weather/airspace feeds; priority policy; audit exports; SSO; dashboards."
], size=20)

prs.save(ppt_path)
print(f"Saved: {ppt_path}")
